"""
Generate End-Semester Viva Presentation (.pptx)
Resume Screening and Gap Analysis – BITS Pilani CCZG628T
Date: 15th May 2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour palette ──────────────────────────────────────────────────────────
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
CARD_BG = RGBColor(0x1A, 0x22, 0x36)
WHITE = RGBColor(0xF1, 0xF5, 0xF9)
GREY = RGBColor(0x94, 0xA3, 0xB8)
MUTED = RGBColor(0x64, 0x74, 0x8B)
CYAN = RGBColor(0x22, 0xD3, 0xEE)
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
ORANGE = RGBColor(0xFB, 0x92, 0x3C)
GREEN = RGBColor(0x34, 0xD3, 0x99)
RED = RGBColor(0xF8, 0x71, 0x71)
YELLOW = RGBColor(0xFB, 0xBF, 0x24)
BLUE = RGBColor(0x60, 0xA5, 0xFA)
PINK = RGBColor(0xF4, 0x72, 0xB6)


def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)

    # Eyebrow
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "BITS Pilani · CCZG628T Dissertation · End-Semester Viva"
    p.font.size = Pt(11)
    p.font.color.rgb = CYAN
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Main title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Resume Screening\nand Gap Analysis"
    p.font.size = Pt(44)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    txBox = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "An NLP-driven pipeline for automated candidate evaluation, intelligent job matching,\ndecision-driven email notifications, and personalised learning path generation."
    p.font.size = Pt(14)
    p.font.color.rgb = GREY
    p.alignment = PP_ALIGN.CENTER

    # Meta info
    meta_lines = [
        "Soumik Basu  |  2024MT03102  |  M.Tech Cloud Computing – WILP",
        "Supervisor: Greta Ghosh  |  Date: 15th May 2026",
        "https://github.com/mikbazz19/bits-dissertation"
    ]
    txBox = slide.shapes.add_textbox(Inches(1), Inches(5.8), Inches(8), Inches(1.2))
    tf = txBox.text_frame
    for i, line in enumerate(meta_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = MUTED
        p.alignment = PP_ALIGN.CENTER
    # Make the GitHub link - inject color directly into run XML (overrides theme hyperlink color)
    from pptx.oxml.ns import qn
    from lxml import etree
    last_p = tf.paragraphs[-1]
    r = last_p.runs[0]
    rPr = r._r.get_or_add_rPr()
    # solidFill must be inserted BEFORE hlinkClick in rPr child order
    solidFill = etree.SubElement(rPr, qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', 'C084FC')
    hlinkClick = rPr.makeelement(qn('a:hlinkClick'), {})
    hlinkClick.set(qn('r:id'), slide.part.relate_to('https://github.com/mikbazz19/bits-dissertation', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink', is_external=True))
    rPr.append(hlinkClick)


def add_section_slide(prs, number, badge_text, title, subtitle, accent=CYAN):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Badge
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{number:02d}  |  {badge_text}"
    p.font.size = Pt(10)
    p.font.color.rgb = accent
    p.font.bold = True

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.font.bold = True

    # Subtitle
    if subtitle:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8.5), Inches(0.6))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(13)
        p.font.color.rgb = GREY

    return slide


def add_bullets(slide, items, left=0.5, top=2.2, width=9.0, height=5.0, font_size=13, color=GREY):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(6)


def add_two_col_bullets(slide, left_items, right_items, top=2.2, left_title="", right_title=""):
    # Left column
    if left_title:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(top - 0.4), Inches(4.3), Inches(0.35))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(11)
        p.font.color.rgb = CYAN
        p.font.bold = True

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(4.3), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(12)
        p.font.color.rgb = GREY
        p.space_after = Pt(4)

    # Right column
    if right_title:
        txBox = slide.shapes.add_textbox(Inches(5.2), Inches(top - 0.4), Inches(4.3), Inches(0.35))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(11)
        p.font.color.rgb = ORANGE
        p.font.bold = True

    txBox = slide.shapes.add_textbox(Inches(5.2), Inches(top), Inches(4.3), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(12)
        p.font.color.rgb = GREY
        p.space_after = Pt(4)


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ═══ SLIDE 1: TITLE ═══
    add_title_slide(prs)

    # ═══ SLIDE 2: MOTIVATION ═══
    slide = add_section_slide(prs, 1, "PROBLEM & MOTIVATION",
        "Why Automated Resume Screening?",
        "Manual screening is slow, inconsistent, and biased — yet it remains the industry norm.",
        CYAN)
    add_bullets(slide, [
        "250+ applications per corporate job posting (LinkedIn/SHRM)",
        "6–7 seconds average recruiter time per resume (Ladders, 2018)",
        "27 million US workers screened out by ATS keyword mismatch (Harvard, 2021)",
        "75% qualified candidates rejected before human review (Ceridian)",
        "",
        "This system addresses: Speed · Consistency · Bidirectional Value",
        "Serves both recruiters (shortlisting) AND candidates (gap report + email)"
    ])

    # ═══ SLIDE 3: OBJECTIVES ═══
    slide = add_section_slide(prs, 2, "OBJECTIVES",
        "11 Core Objectives",
        "Each objective maps directly to a system module or capability.",
        PURPLE)
    add_two_col_bullets(slide, [
        "O1: Multi-format parsing (PDF/DOCX/TXT/OCR)",
        "O2: Information extraction (name, skills, exp, edu)",
        "O3: Job description analysis",
        "O4: Configurable weighted matching + explainability",
        "O5: Gap analysis with learning paths",
        "O6: Interactive Streamlit web interface",
    ], [
        "O7: Batch processing (up to 10 resumes)",
        "O8: Analytics dashboard + quality scoring",
        "O9: Decision-driven email (Accept/Reject/Review)",
        "O10: Unit testing (42+ pytest tests)",
        "O11: Evaluation (P/R/F1 on 10 pairs, Spearman's ρ)",
    ], top=2.2)

    # ═══ SLIDE 4: LITERATURE ═══
    slide = add_section_slide(prs, 3, "REVIEW OF LITERATURE",
        "Prior Work & Research Gap",
        "6 academic papers + 5 commercial tools surveyed. No system covers the full feature set.",
        BLUE)
    add_two_col_bullets(slide, [
        "Faliagka (2012): ML + Semantic Web → domain-specific",
        "Maheshwary & Misra (2018): TF-IDF → no gap analysis",
        "Luo (2019): BERT → expensive, large corpus needed",
        "Roy (2020): spaCy NER → no OCR",
        "Zhang & Wang (2021): GNN → structured data only",
        "Gugnani & Misra (2020): LSTM+CRF → no usable system",
    ], [
        "RESEARCH GAP FILLED:",
        "Open-source & transparent",
        "Multi-modal input (OCR)",
        "Configurable matching + explainability",
        "Batch processing + analytics",
        "Gap analysis + learning paths",
        "Decision-driven email dispatch",
        "Privacy-preserving (100% local)",
    ], top=2.2, left_title="Academic Research", right_title="This System")

    # ═══ SLIDE 5: ARCHITECTURE ═══
    slide = add_section_slide(prs, 4, "SYSTEM ARCHITECTURE",
        "Six-Layer Modular Architecture",
        "Typed interfaces (dataclasses). Each layer communicates only with adjacent layers.",
        CYAN)
    layers = [
        ("UI Layer", "Streamlit (4 Tabs · Single + Batch · Theme · Sidebar)", CYAN),
        ("Data Input Layer", "ResumeParser · OCRProcessor · JDParser · TextPreprocessor", GREEN),
        ("Extraction Layer", "EntityExtractor · SkillExtractor · ExperienceExtractor · Skills DB (118)", ORANGE),
        ("Matching Layer", "JobMatcher (configurable weights) · SimilarityCalculator · Hard Filters", RED),
        ("Analysis & Output", "GapAnalyzer · ReportGenerator · Analytics · EmailSender", PURPLE),
        ("Data & Utilities", "Config · Data Models · Sample Data · Evaluation Package (10 pairs)", BLUE),
    ]
    top = 2.3
    for i, (name, desc, color) in enumerate(layers):
        y = top + i * 0.82
        # Layer name
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(2.2), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(12)
        p.font.color.rgb = color
        p.font.bold = True
        # Layer description
        txBox = slide.shapes.add_textbox(Inches(2.8), Inches(y), Inches(6.8), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = GREY

    # ═══ SLIDE 6: METHODOLOGY ═══
    slide = add_section_slide(prs, 5, "METHODOLOGY",
        "Matching Algorithm & Decision Pipeline",
        "Configurable weighted scoring with hard filters, confidence factor, and 3-outcome classification.",
        ORANGE)
    formula_items = [
        "Overall = w_skills × Skill + w_exp × Experience + w_edu × Education × Confidence",
        "Default weights: Skills 50% · Experience 30% · Education 20%",
        "Skill Score = 0.70 × (matched_req/total_req) + 0.30 × (matched_pref/total_pref)",
        "Experience Score = min(1.0, candidate_years / required_years)",
        "Education Score = 1.0 (match) | 0.5 (any degree) | 0.0 (none)",
        "",
        "Decision Thresholds:",
        "  ≥ 75%  →  ACCEPT (strong match)",
        "  60–74% →  REVIEW (human evaluation needed)",
        "  < 60%  →  REJECT (weak match)",
        "",
        "Hard Filters: min experience | mandatory skills | overqualification (>3 yr)",
        "Fairness: ±5 pts of threshold → advisory for human review",
    ]
    add_bullets(slide, formula_items, font_size=12)

    # ═══ SLIDE 7: KEY MODULES ═══
    slide = add_section_slide(prs, 6, "KEY MODULES",
        "Complete System Implementation",
        "16 modules: parsing, extraction, matching, analysis, email, analytics, evaluation.",
        GREEN)
    add_two_col_bullets(slide, [
        "Data Parsing (PDF/DOCX/TXT/OCR)",
        "Entity Extraction (name, email, phone, edu, certs)",
        "Skill Extraction (118 skills, 8 categories)",
        "Experience Extraction (date parsing, duration)",
        "OCR Processing (Tesseract + OpenCV pipeline)",
        "Matching Engine (weighted + hard filters)",
        "Gap Analyzer (skill/exp/edu + learning path)",
        "Report Generator (Screening/Gap/Comparison PDF)",
    ], [
        "Email Sender (SMTP TLS/SSL, 3 scenarios)",
        "Batch Processing (up to 10 resumes)",
        "Analytics Dashboard (quality score, charts)",
        "Configurable Weights (sidebar sliders)",
        "Explainability Panel ('Why this score?')",
        "Hard Filters + Confidence Factor",
        "Web Application (Streamlit, 4 tabs)",
        "Data Models + Configuration (dataclasses)",
    ], top=2.2)

    # ═══ SLIDE 8: EMAIL MODULE ═══
    slide = add_section_slide(prs, 7, "EMAIL DISPATCH",
        "Decision-Driven Email Notifications",
        "Three scenarios triggered by the matching decision — context-aware communication.",
        PINK)
    add_bullets(slide, [
        "ACCEPT — Notify candidate of selection; facilitate interview scheduling",
        "  → Subject: 'Congratulations! Interview Invitation for [Role]'",
        "  → Body: Selection confirmation + next steps + scheduling link",
        "",
        "REJECT — Communicate outcome with reasons; attach gap analysis report",
        "  → Subject: 'Application Update for [Role]'",
        "  → Body: Polite rejection + specific gaps + attached PDF gap report",
        "  → Value: Helps candidate for professional development",
        "",
        "REVIEW — Escalate borderline candidates to hiring manager",
        "  → Subject: 'Review Required: [Candidate] for [Role]'",
        "  → Body: Summary dossier (score breakdown, matched skills, borderline reason)",
        "",
        "Technical: SMTP TLS (587) / SSL (465) · Gmail/Outlook/Yahoo · MIME attachments",
    ], font_size=11)

    # ═══ SLIDE 9: EVALUATION ═══
    slide = add_section_slide(prs, 8, "EVALUATION",
        "Quantitative Results — 10 Resume–JD Pairs",
        "Expanded from mid-semester's 2 resumes to 10 resume–JD pairs (9 unique annotated resumes).",
        GREEN)
    metrics = [
        "EXTRACTION F1-SCORES (per module):",
        "  Email: 1.00  |  Name: 1.00  |  Phone: 1.00  |  Education: 1.00",
        "  Skill: 0.86  |  Experience: 0.83  |  Certification: 0.75",
        "  MACRO F1 AVERAGE = 0.7981",
        "",
        "MATCHING METRICS:",
        "  Spearman's ρ (rank correlation) = 0.9429",
        "  Shortlisting Accuracy = 81.82% (9/11 correct)",
        "  Mean Absolute Error = 10.56 pts",
        "",
        "WHY F1 AS PRIMARY METRIC?",
        "  Balances precision (no false positives) and recall (no misses)",
        "  Single interpretable number [0–1] for each module",
        "  Standard NLP evaluation practice",
    ]
    add_bullets(slide, metrics, font_size=12)

    # ═══ SLIDE 10: IMPROVEMENT ═══
    slide = add_section_slide(prs, 9, "RESULTS & DISCUSSION",
        "Improvement from Mid-Semester",
        "Significant progress across all dimensions since mid-semester submission.",
        CYAN)
    add_two_col_bullets(slide, [
        "Evaluation pairs: 4 → 10 (+5×)",
        "Experience F1: 0.00 → 0.83 (fixed)",
        "Education F1: 0.00 → 1.00 (fixed)",
        "Spearman's ρ: 1.000 → 0.9429 (robust on larger set)",
        "UI tabs: 3 → 4 (added Analytics)",
        "Processing mode: Single only → Single + Batch",
    ], [
        "POST MID-SEM ADDITIONS:",
        "Batch processing mode (up to 10)",
        "Analytics dashboard + quality scoring",
        "Configurable matching weights (sliders)",
        "Explainability panel",
        "Hard filters + confidence factor",
        "Decision-driven email dispatch",
        "Candidate ranking leaderboard (medals)",
        "CSV export for batch comparison",
    ], top=2.2, left_title="METRIC IMPROVEMENT", right_title="NEW FEATURES")

    # ═══ SLIDE 11: DRAWBACKS ═══
    slide = add_section_slide(prs, 10, "DRAWBACKS & LIMITATIONS",
        "Known Limitations",
        "Honest assessment of current system boundaries.",
        RED)
    add_bullets(slide, [
        "Keyword-based only — no BERT/embedding semantic similarity yet",
        "English language only — multi-language NER requires additional models",
        "10 annotated pairs — meaningful but insufficient for publication-grade significance",
        "Certification extraction F1 = 0.75 — niche formats missed",
        "SMTP firewall issues — corporate networks block ports 587/465",
        "Batch size cap (10) — enterprise scale needs async workers / job queue",
        "No resume generation — only analysis and gap identification",
    ])

    # ═══ SLIDE 12: CONCLUSION ═══
    slide = add_section_slide(prs, 11, "CONCLUSION & FUTURE SCOPE",
        "Summary & Next Steps",
        "All 11 objectives achieved. System fully operational, tested, and evaluated.",
        PURPLE)
    add_two_col_bullets(slide, [
        "KEY CONTRIBUTIONS:",
        "Complete open-source NLP pipeline",
        "Multi-modal input (PDF/DOCX/TXT/OCR)",
        "Configurable matching + explainability",
        "Bidirectional: recruiter + candidate",
        "Decision-driven email notifications",
        "Batch mode + analytics dashboard",
        "F1 = 0.80 | ρ = 0.94 | Acc = 82%",
        "100% local / privacy-preserving",
    ], [
        "FUTURE WORK:",
        "BERT / Sentence-BERT semantic matching",
        "50+ annotated resumes (multi-domain)",
        "Multi-language support (Hindi, French, DE)",
        "OAuth 2.0 for email authentication",
        "Async processing (Celery/Redis)",
        "Resume generation / rewriting suggestions",
        "",
        "11/11 Objectives Achieved ✓",
    ], top=2.2, left_title="DELIVERED", right_title="FUTURE SCOPE")

    # ═══ SLIDE 13: DEMO ═══
    slide = add_section_slide(prs, 12, "LIVE DEMO",
        "Working Application Demo",
        "Demonstrating: parse → match → gap → analytics → email (Streamlit at localhost:8501)",
        CYAN)
    add_bullets(slide, [
        "Step 1: Parse Resume (upload PDF/DOCX/TXT or paste text)",
        "Step 2: Parse Job Description (upload or paste)",
        "Step 3: Calculate Match Score → View decision + explainability",
        "Step 4: Generate Gap Analysis → Download PDF / Send email",
        "Step 5: Switch to Batch Mode → Upload 10 resumes → Leaderboard",
        "Step 6: Analytics tab → Quality score (single) / Batch insights",
        "",
        "Commands:",
        "  streamlit run ui/app.py",
        "  python evaluate.py",
        "  pytest tests/ -v",
    ], font_size=12)

    # ═══ SLIDE 14: THANK YOU ═══
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(48)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Questions & Discussion"
    p.font.size = Pt(20)
    p.font.color.rgb = GREY
    p.alignment = PP_ALIGN.CENTER

    return prs


if __name__ == "__main__":
    print("Generating end-semester viva presentation (PPTX)...")
    prs = build_presentation()
    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                               "dissertation_endsem_viva_presentation.pptx")
    prs.save(output_path)
    print(f"  Saved: {output_path}")
    print(f"  Slides: {len(prs.slides)}")
    print("Done.")
